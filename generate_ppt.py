import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "EventHub — The Future of Event Management"
    subtitle.text = "A Production-Ready Platform for Discovering and Managing Events\nPresented by: [Your Name]"

    # Add the generated image if it exists
    img_path = r"C:\Users\VEERESH AWARALLI\.gemini\antigravity\brain\79a873b2-3e83-43ae-822e-a8c3bdc357b5\eventhub_title_slide_1774761646574.png"
    if os.path.exists(img_path):
        # Add image to the side (half width)
        slide.shapes.add_picture(img_path, Inches(5), Inches(2), width=Inches(4))

    # --- Slide 2: Problem Statement ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Problem Statement"
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "The Challenge in Event Management:"
    p = tf.add_paragraph()
    p.text = "• Fragmentation: Multiple tools for registration, ticketing, and analytics."
    p = tf.add_paragraph()
    p.text = "• Complexity: Setting up custom registration forms is often technical."
    p = tf.add_paragraph()
    p.text = "• Engagement: Lack of built-in social tools like reviews and discussions."
    p = tf.add_paragraph()
    p.text = "• Verification: Manual check-ins are slow and error-prone."

    # --- Slide 3: The Solution ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "The Solution — EventHub"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "An Integrated Ecosystem:"
    p = tf.add_paragraph()
    p.text = "• Single Point of Truth for organizers and attendees."
    p = tf.add_paragraph()
    p.text = "• Automated Workflow: From discovery to QR-code entry."
    p = tf.add_paragraph()
    p.text = "• Data-Driven: Real-time insights for organizers."
    p = tf.add_paragraph()
    p.text = "• AI-Enhanced: Smart discovery via natural language."

    # --- Slide 4: Key Features ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Key Features"
    tf = slide.shapes.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "🎫 QR-Powered Ticketing: Secure entry verification."
    p = tf.add_paragraph()
    p.text = "📈 Organizer Dashboard: Live revenue and registration charts."
    p = tf.add_paragraph()
    p.text = "📝 Custom Forms: Adaptive questions per event."
    p = tf.add_paragraph()
    p.text = "📧 Notifications: Automated registration and reminder emails."
    p = tf.add_paragraph()
    p.text = "💳 Payments: Razorpay and 'Pay at Venue' support."

    # --- Slide 5: System Architecture ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "System Architecture"
    tf = slide.shapes.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "• Backend Core: Django 5.1 & Python 3.12"
    p = tf.add_paragraph()
    p.text = "• Database Layer: PostgreSQL for robust data"
    p = tf.add_paragraph()
    p.text = "• Media Storage: Cloudinary integration"
    p = tf.add_paragraph()
    p.text = "• Frontend: Clean design with Bootstrap 5"
    p = tf.add_paragraph()
    p.text = "• AI Integration: Google Gemini API Assistant"

    # --- Slide 6: Flow Diagram - User Onboarding ---
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only
    slide.shapes.title.text = "User Onboarding Flow"
    # Note: Complex diagrams are best added as screenshots of the Mermaid output 
    # included in presentation_content.md.
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tf = subtitle.text_frame
    tf.text = "[Insert Flow Diagram A from presentation_content.md here]"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # --- Slide 7: Flow Diagram - Booking Flow ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Event Booking & Ticketing Flow"
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tf = subtitle.text_frame
    tf.text = "[Insert Flow Diagram B from presentation_content.md here]"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # --- Slide 8: Organizer Experience ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Organizer Experience"
    tf = slide.shapes.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "• Revenue Overview: Visual breakdown of ticket sales."
    p = tf.add_paragraph()
    p.text = "• Attendee Insights: Exportable CSV lists."
    p = tf.add_paragraph()
    p.text = "• Performance Tracking: Analytics at a glance."
    p = tf.add_paragraph()
    p.text = "• Effortless Management: Complete CRUD for events."

    # --- Slide 9: AI Discovery Engine ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "AI Discovery Engine"
    tf = slide.shapes.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "Concept: Hybrid Assistant (DB Queries + AI Reasoning)"
    p = tf.add_paragraph()
    p.text = "• Natural Language Discovery ('What's on this weekend?')"
    p = tf.add_paragraph()
    p.text = "• Process Guidance: Helping users with registration."
    p = tf.add_paragraph()
    p.text = "• Context Awareness: Location-based responses (Kalaburagi)."

    # --- Slide 10: Conclusion & Roadmap ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Conclusion & Roadmap"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Current State: Stable, Production-Ready."
    p = tf.add_paragraph()
    p.text = "Future Goals:"
    p = tf.add_paragraph()
    p.text = "• Google Calendar Integration"
    p = tf.add_paragraph()
    p.text = "• Native QR Scanner Mobile Apps"
    p = tf.add_paragraph()
    p.text = "• Multi-lingual support"

    prs.save('EventHub_Project_Presentation.pptx')
    print("Presentation created successfully as 'EventHub_Project_Presentation.pptx'")

if __name__ == "__main__":
    create_presentation()
